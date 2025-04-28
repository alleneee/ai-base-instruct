"""文档处理器模块"""
import os
import logging
import uuid
import re
import mimetypes
from datetime import datetime
from typing import Dict, List, Optional, Any, Union, BinaryIO, Tuple
from pathlib import Path
import json
import math

from llama_index.core import Document, Settings
from llama_index.core.ingestion import IngestionPipeline
from llama_index.core.node_parser import SentenceSplitter, MarkdownNodeParser, HTMLNodeParser, CodeSplitter
from llama_index.readers.file import PyMuPDFReader, DocxReader, UnstructuredReader
from llama_index.embeddings.openai import OpenAIEmbedding
from llama_index.vector_stores.milvus import MilvusVectorStore, IndexManagement
from markitdown import MarkItDown
from llama_index.core.callbacks import CallbackManager, TokenCountingHandler
import tiktoken

from enterprise_kb.core.config.settings import settings
from enterprise_kb.storage.vector_store import get_storage_context
from enterprise_kb.utils.milvus_client import get_milvus_client
from enterprise_kb.storage.datasource.milvus import MilvusDataSource, MilvusDataSourceConfig

logger = logging.getLogger(__name__)

class DocumentComplexity:
    """文档复杂度评估结果"""
    LOW = "low"        # 简单文本，适合直接处理
    MEDIUM = "medium"  # 中等复杂度，标准处理
    HIGH = "high"      # 高复杂度，需要特殊处理

class DocumentFeatures:
    """文档特征"""
    def __init__(self):
        self.page_count: int = 0            # 页数
        self.has_tables: bool = False       # 是否包含表格
        self.has_images: bool = False       # 是否包含图片
        self.has_code: bool = False         # 是否包含代码
        self.text_density: float = 0.0      # 文本密度
        self.estimated_tokens: int = 0      # 估计的令牌数
        self.language: str = "unknown"      # 文档语言
        self.structure_level: int = 0       # 结构层级数量（标题等）
        self.avg_sentence_length: float = 0 # 平均句子长度
        self.metadata: Dict[str, Any] = {}  # 其他元数据

    def to_dict(self) -> Dict[str, Any]:
        """转换为字典"""
        return {
            "page_count": self.page_count,
            "has_tables": self.has_tables,
            "has_images": self.has_images,
            "has_code": self.has_code,
            "text_density": self.text_density,
            "estimated_tokens": self.estimated_tokens,
            "language": self.language,
            "structure_level": self.structure_level,
            "avg_sentence_length": self.avg_sentence_length,
            "metadata": self.metadata
        }

class DocumentAnalyzer:
    """文档分析器，用于评估文档复杂度和特征"""

    def __init__(self):
        """初始化文档分析器"""
        # mime类型初始化
        mimetypes.init()

    def analyze_document(self, file_path: str) -> Tuple[DocumentFeatures, DocumentComplexity]:
        """分析文档并返回特征和复杂度评估

        Args:
            file_path: 文档路径

        Returns:
            文档特征和复杂度评估
        """
        file_ext = os.path.splitext(file_path)[1].lower()
        features = DocumentFeatures()

        # 根据文件类型选择适当的分析方法
        if file_ext in ['.pdf']:
            self._analyze_pdf(file_path, features)
        elif file_ext in ['.docx', '.doc']:
            self._analyze_word(file_path, features)
        elif file_ext in ['.md', '.markdown']:
            self._analyze_markdown(file_path, features)
        elif file_ext in ['.html', '.htm']:
            self._analyze_html(file_path, features)
        elif file_ext in ['.py', '.js', '.java', '.cpp', '.c', '.rb']:
            self._analyze_code(file_path, features)
        elif file_ext in ['.csv', '.xlsx', '.xls']:
            self._analyze_table(file_path, features)
        elif file_ext in ['.txt']:
            self._analyze_text(file_path, features)
        elif file_ext in ['.ppt', '.pptx']:
            self._analyze_presentation(file_path, features)
        else:
            self._analyze_generic(file_path, features)

        # 根据特征判断复杂度
        complexity = self._determine_complexity(features)

        return features, complexity

    def _analyze_pdf(self, file_path: str, features: DocumentFeatures):
        """分析PDF文档"""
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(file_path)
            features.page_count = len(doc)

            # 统计图片数
            image_count = 0
            table_count = 0
            text_blocks = 0
            total_text_length = 0

            # 分析前10页或全部页面（取较小值）
            analyze_pages = min(10, features.page_count)
            for page_idx in range(analyze_pages):
                page = doc[page_idx]

                # 文本分析
                text = page.get_text()
                if text:
                    text_blocks += 1
                    total_text_length += len(text)

                # 提取图片
                image_list = page.get_images(full=True)
                image_count += len(image_list)

                # 简单表格检测（基于文本模式）
                if '|' in text or '\t' in text:
                    table_count += 1

            # 设置特征
            features.has_images = image_count > 0
            features.has_tables = table_count > 0

            # 估算文本密度和结构
            if analyze_pages > 0:
                features.text_density = total_text_length / analyze_pages if analyze_pages else 0

            # 估算token数
            features.estimated_tokens = total_text_length // 4  # 粗略估计：平均每个token 4个字符

            # 检测语言
            features.language = self._detect_language(doc[0].get_text(100))  # 使用第一页的前100个字符

            # 结构分析
            toc = doc.get_toc()
            features.structure_level = len(set(lvl for lvl, _, _ in toc)) if toc else 0

        except Exception as e:
            logger.error(f"PDF分析失败: {str(e)}")

    def _analyze_word(self, file_path: str, features: DocumentFeatures):
        """分析Word文档"""
        try:
            from docx import Document as DocxDocument

            doc = DocxDocument(file_path)

            # 分析段落数和表格
            features.page_count = len(doc.paragraphs) // 40  # 粗略估计：每页40段
            features.has_tables = len(doc.tables) > 0

            # 分析图片（间接判断）
            try:
                for rel in doc.part.rels.values():
                    if "image" in rel.target_ref:
                        features.has_images = True
                        break
            except:
                pass

            # 文本分析
            text = "\n".join(p.text for p in doc.paragraphs)
            features.estimated_tokens = len(text) // 4

            # 语言检测
            if text:
                features.language = self._detect_language(text[:100])

            # 简单结构分析
            headings = 0
            for p in doc.paragraphs:
                if p.style and "Heading" in p.style.name:
                    headings += 1

            features.structure_level = 3 if headings > 10 else (2 if headings > 3 else 1)

        except Exception as e:
            logger.error(f"Word文档分析失败: {str(e)}")

    def _analyze_markdown(self, file_path: str, features: DocumentFeatures):
        """分析Markdown文档"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()

            # 分析标题等级
            heading_pattern = re.compile(r'^(#{1,6})\s+', re.MULTILINE)
            headings = heading_pattern.findall(content)
            heading_levels = [len(h) for h in headings]
            features.structure_level = len(set(heading_levels)) if heading_levels else 0

            # 分析代码块
            code_blocks = re.findall(r'```[\s\S]*?```', content)
            features.has_code = len(code_blocks) > 0

            # 分析表格
            features.has_tables = '|---' in content or '| ---' in content

            # 分析图片
            features.has_images = '![' in content

            # 文本分析
            features.page_count = content.count('\n') // 40 + 1  # 粗略估计
            features.estimated_tokens = len(content) // 4

            # 语言检测
            features.language = self._detect_language(content[:100])

        except Exception as e:
            logger.error(f"Markdown分析失败: {str(e)}")

    def _analyze_html(self, file_path: str, features: DocumentFeatures):
        """分析HTML文档"""
        try:
            from bs4 import BeautifulSoup

            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()

            soup = BeautifulSoup(content, 'html.parser')

            # 基本特征分析
            features.has_tables = len(soup.find_all('table')) > 0
            features.has_images = len(soup.find_all('img')) > 0
            features.has_code = len(soup.find_all(['code', 'pre'])) > 0

            # 文本分析
            text = soup.get_text()
            features.estimated_tokens = len(text) // 4
            features.page_count = len(text) // 3000 + 1  # 粗略估计

            # 结构分析
            headings = soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6'])
            features.structure_level = len(set(h.name for h in headings)) if headings else 0

            # 语言检测
            features.language = self._detect_language(text[:100])

        except Exception as e:
            logger.error(f"HTML分析失败: {str(e)}")

    def _analyze_code(self, file_path: str, features: DocumentFeatures):
        """分析代码文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()

            features.has_code = True
            features.page_count = content.count('\n') // 50 + 1
            features.estimated_tokens = len(content) // 3  # 代码token比自然语言短

            # 检测注释占比、类/函数定义等
            # 这里只是简化实现，实际上可以针对不同编程语言有更专门的分析
            comment_lines = 0
            if file_path.endswith('.py'):
                comment_lines = content.count('#') + content.count('"""') // 2
            elif file_path.endswith(('.js', '.java', '.cpp', '.c')):
                comment_lines = content.count('//') + content.count('/*') + content.count('*/')

            features.metadata['comment_ratio'] = comment_lines / (content.count('\n') + 1) if content.count('\n') > 0 else 0
            features.language = os.path.splitext(file_path)[1][1:]  # 使用文件扩展名作为语言

        except Exception as e:
            logger.error(f"代码分析失败: {str(e)}")

    def _analyze_table(self, file_path: str, features: DocumentFeatures):
        """分析表格文件(CSV/Excel)"""
        try:
            features.has_tables = True

            if file_path.endswith('.csv'):
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()

                rows = content.count('\n') + 1
                features.page_count = max(1, rows // 100)  # 粗略估计：每页100行
                features.estimated_tokens = len(content) // 5

                # 分析列数
                if '\n' in content:
                    first_line = content.split('\n')[0]
                    features.metadata['columns'] = first_line.count(',') + 1

            elif file_path.endswith(('.xlsx', '.xls')):
                try:
                    import pandas as pd
                    df = pd.read_excel(file_path)
                    features.page_count = max(1, len(df) // 100)
                    features.estimated_tokens = df.size * 2  # 每个单元格约2个token
                    features.metadata['columns'] = len(df.columns)
                    features.metadata['rows'] = len(df)
                except ImportError:
                    # 如果pandas未安装，使用更基本的方法
                    features.page_count = 1

        except Exception as e:
            logger.error(f"表格分析失败: {str(e)}")

    def _analyze_text(self, file_path: str, features: DocumentFeatures):
        """分析纯文本文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()

            features.page_count = content.count('\n') // 40 + 1
            features.estimated_tokens = len(content) // 4

            # 检测是否包含表格式文本
            lines = content.split('\n')
            if len(lines) > 3:
                table_like_lines = sum(1 for line in lines if line.count('\t') > 2 or line.count('|') > 2)
                features.has_tables = table_like_lines > 3

            # 语言检测
            features.language = self._detect_language(content[:100])

            # 简单分析句子长度
            sentences = re.split(r'[.!?]+', content)
            if sentences:
                features.avg_sentence_length = sum(len(s) for s in sentences) / len(sentences)

        except Exception as e:
            logger.error(f"文本分析失败: {str(e)}")

    def _analyze_presentation(self, file_path: str, features: DocumentFeatures):
        """分析演示文稿"""
        try:
            import pptx
            presentation = pptx.Presentation(file_path)

            features.page_count = len(presentation.slides)

            text_content = []
            image_count = 0
            table_count = 0

            # 分析每个幻灯片
            for slide in presentation.slides:
                # 提取文本
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)

                    # 检测图片
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        image_count += 1

                    # 检测表格
                    if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                        table_count += 1

            # 设置特征
            features.has_images = image_count > 0
            features.has_tables = table_count > 0
            features.estimated_tokens = sum(len(text) for text in text_content) // 4

            # 语言检测
            if text_content:
                sample_text = text_content[0][:100] if text_content[0] else ""
                features.language = self._detect_language(sample_text)

        except Exception as e:
            logger.error(f"演示文稿分析失败: {str(e)}")
            # 基本估计
            features.page_count = 1
            features.has_images = True

    def _analyze_generic(self, file_path: str, features: DocumentFeatures):
        """通用文件分析"""
        # 尝试判断MIME类型
        mime_type, _ = mimetypes.guess_type(file_path)

        # 设置基本特征
        features.page_count = 1

        if mime_type:
            features.metadata['mime_type'] = mime_type

            if 'image' in mime_type:
                features.has_images = True
            elif 'text' in mime_type:
                self._analyze_text(file_path, features)

        # 文件大小作为复杂度参考
        try:
            file_size = os.path.getsize(file_path)
            features.metadata['file_size'] = file_size
            features.estimated_tokens = file_size // 10  # 非常粗略的估计
        except:
            pass

    def _detect_language(self, text_sample: str) -> str:
        """检测文本语言

        这是一个非常简化的实现，实际应用中可以使用更复杂的语言检测库如langdetect
        """
        if not text_sample:
            return "unknown"

        # 简化的语言检测逻辑
        chinese_chars = len(re.findall(r'[\u4e00-\u9fff]', text_sample))
        if chinese_chars > len(text_sample) * 0.2:
            return "zh"

        # 可以添加其他语言的检测
        # 默认返回英文
        return "en"

    def _determine_complexity(self, features: DocumentFeatures) -> DocumentComplexity:
        """根据特征确定文档复杂度"""
        # 高复杂度指标
        if features.page_count > 50 or features.estimated_tokens > 50000:
            return DocumentComplexity.HIGH

        # 包含复杂元素且规模较大
        if (features.has_tables or features.has_images) and features.page_count > 20:
            return DocumentComplexity.HIGH

        # 中等复杂度指标
        if features.page_count > 10 or features.estimated_tokens > 10000:
            return DocumentComplexity.MEDIUM

        if features.has_tables or features.has_images or features.has_code:
            return DocumentComplexity.MEDIUM

        # 其他为低复杂度
        return DocumentComplexity.LOW

class AdaptiveChunkingStrategy:
    """自适应分块策略"""

    def __init__(self):
        """初始化自适应分块策略"""
        self.default_chunk_size = settings.CHUNK_SIZE
        self.default_chunk_overlap = settings.CHUNK_OVERLAP

    def get_chunking_parameters(
        self,
        doc_type: str,
        features: DocumentFeatures,
        complexity: DocumentComplexity
    ) -> Dict[str, Any]:
        """获取分块参数

        Args:
            doc_type: 文档类型
            features: 文档特征
            complexity: 文档复杂度

        Returns:
            分块参数
        """
        # 根据文档类型优化
        if doc_type == "markdown":
            return self._get_markdown_parameters(features, complexity)
        elif doc_type == "code":
            return self._get_code_parameters(features)
        elif doc_type == "table":
            return self._get_table_parameters(features)
        elif doc_type == "pdf":
            return self._get_pdf_parameters(features, complexity)
        elif doc_type == "text":
            return self._get_text_parameters(features, complexity)
        else:
            # 通用参数
            return self._get_generic_parameters(features, complexity)

    def _get_markdown_parameters(self, features: DocumentFeatures, complexity: DocumentComplexity) -> Dict[str, Any]:
        """获取Markdown文档的分块参数"""
        # Markdown文档可以使用专门的解析器
        params = {
            "chunking_type": "markdown",
            "parser_type": "markdown",
            "chunk_size": self.default_chunk_size,
            "chunk_overlap": self.default_chunk_overlap
        }

        # 根据结构复杂度调整
        if features.structure_level > 3:
            # 有丰富的标题层级，可以适当增加块大小
            params["chunk_size"] = min(2048, self.default_chunk_size * 1.5)

        # 复杂度调整
        if complexity == DocumentComplexity.HIGH:
            # 更小的块以处理复杂内容
            params["chunk_size"] = max(512, params["chunk_size"] // 2)
            params["chunk_overlap"] = min(150, params["chunk_overlap"] * 2)
        elif complexity == DocumentComplexity.LOW:
            # 较大的块以保持上下文连贯性
            params["chunk_size"] = min(2048, params["chunk_size"] * 1.5)

        return params

    def _get_code_parameters(self, features: DocumentFeatures) -> Dict[str, Any]:
        """获取代码文件的分块参数"""
        return {
            "chunking_type": "code",
            "parser_type": "code",
            "language": features.language,
            "chunk_size": min(1024, self.default_chunk_size),
            "chunk_overlap": min(100, self.default_chunk_overlap),
            "max_lines": 100
        }

    def _get_table_parameters(self, features: DocumentFeatures) -> Dict[str, Any]:
        """获取表格文件的分块参数"""
        # 表格文件可能需要特殊处理
        return {
            "chunking_type": "table",
            "parser_type": "table",
            "chunk_size": min(1024, self.default_chunk_size),
            "chunk_overlap": 0,  # 表格数据通常不需要重叠
            "metadata": {
                "columns": features.metadata.get("columns", 0),
                "rows": features.metadata.get("rows", 0)
            }
        }

    def _get_pdf_parameters(self, features: DocumentFeatures, complexity: DocumentComplexity) -> Dict[str, Any]:
        """获取PDF文档的分块参数"""
        params = {
            "chunking_type": "pdf",
            "parser_type": "sentence",
            "chunk_size": self.default_chunk_size,
            "chunk_overlap": self.default_chunk_overlap
        }

        # 有图表的PDF可能需要更细粒度的分块
        if features.has_images or features.has_tables:
            params["chunk_size"] = max(512, params["chunk_size"] // 2)
            params["chunk_overlap"] = min(200, params["chunk_overlap"] * 1.5)

        # 根据复杂度调整
        if complexity == DocumentComplexity.HIGH:
            params["chunk_size"] = max(256, params["chunk_size"] // 2)
        elif complexity == DocumentComplexity.LOW:
            params["chunk_size"] = min(2048, params["chunk_size"] * 1.2)

        return params

    def _get_text_parameters(self, features: DocumentFeatures, complexity: DocumentComplexity) -> Dict[str, Any]:
        """获取文本文件的分块参数"""
        params = {
            "chunking_type": "text",
            "parser_type": "sentence",
            "chunk_size": self.default_chunk_size,
            "chunk_overlap": self.default_chunk_overlap
        }

        # 根据复杂度调整
        if complexity == DocumentComplexity.HIGH:
            params["chunk_size"] = max(512, params["chunk_size"] // 1.5)
        elif complexity == DocumentComplexity.LOW:
            params["chunk_size"] = min(2048, params["chunk_size"] * 1.2)

        # 考虑句子长度
        if features.avg_sentence_length > 100:
            # 长句子需要更大的重叠
            params["chunk_overlap"] = min(200, params["chunk_overlap"] * 1.5)

        return params

    def _get_generic_parameters(self, features: DocumentFeatures, complexity: DocumentComplexity) -> Dict[str, Any]:
        """获取通用参数"""
        params = {
            "chunking_type": "generic",
            "parser_type": "sentence",
            "chunk_size": self.default_chunk_size,
            "chunk_overlap": self.default_chunk_overlap
        }

        # 根据复杂度调整
        if complexity == DocumentComplexity.HIGH:
            params["chunk_size"] = max(512, params["chunk_size"] // 1.5)
            params["chunk_overlap"] = min(200, params["chunk_overlap"] * 1.5)
        elif complexity == DocumentComplexity.LOW:
            params["chunk_size"] = min(2048, params["chunk_size"] * 1.2)

        return params

    def create_parser(self, params: Dict[str, Any]):
        """创建解析器

        根据参数创建适当的分块解析器

        Args:
            params: 分块参数

        Returns:
            节点解析器
        """
        parser_type = params.get("parser_type", "sentence")

        if parser_type == "markdown":
            return MarkdownNodeParser(
                chunk_size=params["chunk_size"],
                chunk_overlap=params["chunk_overlap"]
            )
        elif parser_type == "code":
            return CodeSplitter(
                language=params.get("language", ""),
                chunk_size=params["chunk_size"],
                chunk_overlap=params["chunk_overlap"],
                max_lines=params.get("max_lines", 100)
            )
        elif parser_type == "html":
            return HTMLNodeParser(
                chunk_size=params["chunk_size"],
                chunk_overlap=params["chunk_overlap"]
            )
        else:
            # 默认使用句子分块器
            return SentenceSplitter(
                chunk_size=params["chunk_size"],
                chunk_overlap=params["chunk_overlap"]
            )

class DocumentProcessor:
    """文档处理器"""

    def __init__(self):
        """初始化文档处理器"""
        # 设置LlamaIndex的默认设置
        Settings.embed_model = OpenAIEmbedding()
        Settings.chunk_size = settings.CHUNK_SIZE
        Settings.chunk_overlap = settings.CHUNK_OVERLAP

        # 初始化MarkItDown转换器
        self.markitdown = MarkItDown(enable_plugins=True)

        # 初始化计数器
        self.token_counter = TokenCountingHandler(
            tokenizer=tiktoken.encoding_for_model("gpt-3.5-turbo").encode
        )

        # 初始化存储上下文
        self.storage_context = get_storage_context()

        # 初始化Milvus数据源
        self.milvus_config = MilvusDataSourceConfig(
            name="primary",
            description="主要向量存储",
            uri=f"http://{settings.MILVUS_HOST}:{settings.MILVUS_PORT}",
            collection_name=settings.MILVUS_COLLECTION,
            dimension=settings.MILVUS_DIMENSION
        )
        self.milvus_datasource = MilvusDataSource(self.milvus_config)

    async def process_document(
        self,
        file_path: str,
        metadata: Dict[str, Any],
        datasource_name: Optional[str] = "primary",
        use_markitdown: bool = True,
        use_semantic_chunking: bool = False,
        chunking_strategy: Optional[str] = None,
        store_metadata: bool = False
    ) -> Dict[str, Any]:
        """
        处理文档并添加到向量存储

        Args:
            file_path: 文件路径
            metadata: 文档元数据，应包含doc_id
            datasource_name: 数据源名称
            use_markitdown: 是否使用MarkItDown进行预处理
            use_semantic_chunking: 是否使用语义分块
            chunking_strategy: 分块策略，可选值：sentence, paragraph, fixed, hybrid
            store_metadata: 是否将元数据存储到MySQL数据库

        Returns:
            处理结果
        """
        try:
            # 确保元数据包含doc_id
            if "doc_id" not in metadata:
                metadata["doc_id"] = str(uuid.uuid4())

            doc_id = metadata["doc_id"]
            logger.info(f"开始处理文档: {file_path}, doc_id: {doc_id}")

            # 确定文档类型
            file_ext = os.path.splitext(file_path)[1].lower()
            doc_type = file_ext.lstrip(".")

            # 首先分析文档特征
            features, complexity = await self._analyze_document(file_path, doc_type)

            # 根据文档特征决定是否需要转换为Markdown
            convert_to_markdown = use_markitdown and self._should_convert_to_markdown(doc_type, features, complexity)

            # 使用适当的Reader读取文档
            if convert_to_markdown:
                # 使用MarkItDown转换文档
                logger.info(f"转换文档为Markdown: {file_path}")
                content = self.convert_to_markdown(file_path)
                documents = [Document(text=content, metadata=metadata)]
            else:
                # 使用标准Reader
                reader = self.get_reader_for_file(file_path)
                if reader:
                    documents = reader.load(file_path, metadata=metadata)
                    logger.info(f"使用 {reader.__class__.__name__} 加载文档: {file_path}")
                else:
                    # 回退到纯文本
                    with open(file_path, "r", encoding="utf-8") as f:
                        content = f.read()
                    documents = [Document(text=content, metadata=metadata)]
                    logger.info(f"以纯文本加载文档: {file_path}")

            # 根据文档特征和复杂度选择分块策略
            chunker = self._get_chunker(doc_type, features, complexity, chunking_strategy, use_semantic_chunking)

            # 创建并执行处理管道
            callback_manager = CallbackManager([self.token_counter])
            pipeline = IngestionPipeline(
                transformations=[chunker],
                callback_manager=callback_manager
            )

            # 处理文档
            nodes = pipeline.run(documents=documents)

            # 添加到向量存储
            if nodes:
                # 连接到Milvus数据源
                await self.milvus_datasource.connect()

                try:
                    # 将节点添加到Milvus向量存储
                    logger.info(f"正在将 {len(nodes)} 个节点添加到Milvus向量存储...")

                    # 确保每个节点的元数据中包含doc_id
                    for node in nodes:
                        if node.metadata is None:
                            node.metadata = {}
                        node.metadata["doc_id"] = doc_id
                        node.metadata["datasource"] = datasource_name

                    # 添加节点到Milvus
                    node_ids = await self.milvus_datasource.add_documents(nodes)

                    logger.info(f"成功将节点添加到Milvus向量存储，返回的节点ID: {node_ids[:5]}...")
                except Exception as e:
                    logger.error(f"添加节点到Milvus失败: {str(e)}")
                    # 如果添加失败，仍然返回节点ID
                    node_ids = [node.node_id for node in nodes]
                finally:
                    # 断开Milvus连接
                    await self.milvus_datasource.disconnect()

                token_count = self.token_counter.total_embedding_token_count
                logger.info(f"文档处理完成，生成了 {len(nodes)} 个节点，总计使用了 {token_count} 个token")

                # 处理结果
                result = {
                    "doc_id": doc_id,
                    "node_count": len(nodes),
                    "node_ids": node_ids,
                    "token_count": token_count,
                    "status": "completed",
                    "features": features.to_dict(),
                    "complexity": complexity,
                    "datasource": datasource_name,
                    "vector_store": "milvus"  # 添加向量存储类型信息
                }

                # 如果需要存储元数据到MySQL
                if store_metadata:
                    try:
                        # 导入必要的模块
                        from enterprise_kb.db.repositories.document_repository import DocumentRepository
                        from enterprise_kb.schemas.documents import DocumentStatus

                        # 获取文件名
                        file_name = os.path.basename(file_path)

                        # 获取文件大小
                        file_size = os.path.getsize(file_path) if os.path.exists(file_path) else 0

                        # 准备数据库记录
                        document_data = {
                            "id": doc_id,
                            "file_name": file_name,
                            "file_path": file_path,
                            "file_type": doc_type,
                            "title": metadata.get("title", file_name),
                            "description": metadata.get("description"),
                            "status": DocumentStatus.COMPLETED.value,
                            "size_bytes": file_size,
                            "node_count": len(nodes),
                            "metadata": metadata,
                            "created_at": datetime.now(),
                            "updated_at": datetime.now()
                        }

                        # 创建文档仓库
                        doc_repo = DocumentRepository()

                        # 异步存储元数据
                        import asyncio
                        asyncio.create_task(self._store_metadata_to_db(doc_repo, document_data))

                        logger.info(f"已创建异步任务存储文档元数据: {doc_id}")
                    except Exception as e:
                        logger.error(f"准备存储元数据失败: {str(e)}")

                return result
            else:
                logger.warning(f"文档处理未生成节点: {file_path}")
                return {
                    "doc_id": doc_id,
                    "node_count": 0,
                    "status": "warning",
                    "message": "文档处理未生成节点"
                }

        except Exception as e:
            logger.error(f"文档处理失败: {str(e)}", exc_info=True)
            return {
                "doc_id": metadata.get("doc_id"),
                "status": "error",
                "error": str(e)
            }

    async def _store_metadata_to_db(self, doc_repo, document_data: Dict[str, Any]) -> None:
        """
        异步存储元数据到数据库

        Args:
            doc_repo: 文档仓库实例
            document_data: 文档数据
        """
        try:
            # 检查文档是否已存在
            existing_doc = await doc_repo.get(document_data["id"])

            if existing_doc:
                # 更新现有文档
                await doc_repo.update(document_data["id"], document_data)
                logger.info(f"更新文档元数据成功: {document_data['id']}")
            else:
                # 创建新文档
                await doc_repo.create(document_data)
                logger.info(f"创建文档元数据成功: {document_data['id']}")

        except Exception as e:
            logger.error(f"存储元数据到数据库失败: {str(e)}")

    async def delete_document_vectors(self, doc_id: str) -> bool:
        """
        从Milvus向量数据库中删除文档的向量数据

        Args:
            doc_id: 文档ID

        Returns:
            是否成功删除
        """
        try:
            # 连接到Milvus数据源
            await self.milvus_datasource.connect()

            # 删除文档向量
            success = await self.milvus_datasource.delete_document(doc_id)

            # 断开连接
            await self.milvus_datasource.disconnect()

            if success:
                logger.info(f"成功从Milvus删除文档向量: {doc_id}")
            else:
                logger.warning(f"从Milvus删除文档向量可能未成功: {doc_id}")

            return success
        except Exception as e:
            logger.error(f"从Milvus删除文档向量失败: {str(e)}")
            return False

    async def delete_document(self, doc_id: str) -> bool:
        """
        删除文档，包括从Milvus向量数据库中删除向量数据

        Args:
            doc_id: 文档ID

        Returns:
            是否成功删除
        """
        try:
            logger.info(f"开始删除文档: {doc_id}")

            # 1. 从Milvus删除向量数据
            vector_deleted = await self.delete_document_vectors(doc_id)
            if not vector_deleted:
                logger.warning(f"从Milvus删除向量数据失败: {doc_id}")

            # 2. 从数据库删除元数据
            # 这里可以添加从数据库删除元数据的代码
            # 例如：await document_repository.delete(doc_id)

            # 3. 删除相关文件
            # 这里可以添加删除文件的代码
            # 例如：os.remove(file_path) 如果有文件路径的话

            logger.info(f"文档删除完成: {doc_id}")
            return True
        except Exception as e:
            logger.error(f"删除文档失败: {str(e)}")
            raise

    def convert_to_markdown(self, file_path: str) -> str:
        """使用MarkItDown将文档转换为Markdown格式

        Args:
            file_path: 文件路径

        Returns:
            Markdown格式的文本

        Raises:
            RuntimeError: 转换失败
        """
        try:
            # 使用MarkItDown转换文件
            result = self.markitdown.convert(file_path)

            # 如果转换失败，MarkItDown会返回空字符串或None
            if not result:
                # 尝试备用方法
                logger.warning(f"MarkItDown转换返回空结果，尝试使用备用方法: {file_path}")
                # 使用备用方法，例如直接读取文件内容
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    result = f.read()

            # 如果仍然失败，抛出异常
            if not result:
                raise RuntimeError(f"文档转换失败: {file_path}")

            logger.info(f"文档成功转换为Markdown: {file_path}")
            return result

        except Exception as e:
            logger.error(f"Markdown转换失败: {str(e)}")
            raise RuntimeError(f"文档转换失败: {str(e)}")

# 创建单例实例
_document_processor = None

def get_document_processor() -> DocumentProcessor:
    """获取文档处理器单例"""
    global _document_processor
    if _document_processor is None:
        _document_processor = DocumentProcessor()
    return _document_processor